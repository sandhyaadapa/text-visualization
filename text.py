import streamlit as st
from pptx import Presentation
from pptx.util import Inches
import io

st.title("📊 PowerPoint Generator")

st.write("Enter your slide content below:")

# Input form
with st.form("ppt_form"):
    slide_title = st.text_input("Slide Title", "My Presentation Slide")
    bullet_points = st.text_area("Bullet Points (one per line)", "First point\nSecond point\nThird point")
    submit = st.form_submit_button("Generate PPT")

if submit:
    prs = Presentation()
    slide_layout = prs.slide_layouts[1]  # Title and Content layout
    slide = prs.slides.add_slide(slide_layout)

    title = slide.shapes.title
    content = slide.placeholders[1]

    title.text = slide_title
    for line in bullet_points.strip().split("\n"):
        content.text += f"\n• {line.strip()}"

    # Save to BytesIO
    pptx_io = io.BytesIO()
    prs.save(pptx_io)
    pptx_io.seek(0)

    st.success("✅ PowerPoint file generated!")
    st.download_button(
        label="📥 Download PPTX",
        data=pptx_io,
        file_name="generated_presentation.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )
